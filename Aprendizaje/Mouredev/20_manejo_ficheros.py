### Manejo de ficheros - File Handling ###

import json
import os

# 1. .txt file
'''
# Abrir el fichero
# Primero buscará el fichero en la carpeta donde se encuentra el script
txt_file = open("Aprendizaje/Mouredev/mi_fichero.txt", "r+")

# Al no encontrarlo, creará el fichero.
txt_file.write("Mi nombre es Xavi\nMi apellido es Quesada\nTengo 31 años\nY mi lenguaje favorito es Python")

# Leer el fichero
#print(txt_file.read()) #Lee todo el fichero
print(txt_file.read(10)) #Lee los primeros 10 caracteres
print(txt_file.readline()) #Lee la primera línea
print(txt_file.readline()) #Lee la siguiente línea
print(txt_file.readlines()) #Lee todas las líneas y las devuelve en una lista

for line in txt_file.readlines(): #Llamamos a readlines y las junta a un listado y lo recorremos con un for y las lee linea a linea a linea y las imprime
    print(line)

# Escribe en el fichero con un salto de linea y lo guarda en el fichero
txt_file.write("\nSoy de Sabadell")
print(txt_file.readline()) # Lee la línea añadida

# Cierra el fichero
txt_file.close()

with open("Aprendizaje/Mouredev/mi_fichero.txt", "a") as mi_documento: #Con with no hace falta cerrar el fichero
    mi_documento.write("\nTambién me gusta Java") #No se puede escribir en un fichero que se ha abierto con with

# Elimina el fichero
# os.remove("Aprendizaje/Mouredev/mi_fichero.txt")
'''

# 2. .json file
# Json es un formato de datos que se utiliza para intercambiar información entre sistemas
'''
# Crea un fichero json
json_file = open("Aprendizaje/Mouredev/mi_fichero.json", "w+")

# Se crea como un diccionario
json_test = {
    "nombre": "Xavi",
    "apellido": "Quesada",
    "edad": 31,
    "lenguajes": ["Python", "Java", "C++"],
    "ciudad": "Sabadell"
}

# Escribe el diccionario en el fichero
json.dump(json_test, json_file, indent=2) #indent=2 es para que se vea mejor el json haciendo un salto de linea

json_file.close() #Cierra el fichero

# Lee el fichero json
with open("Aprendizaje/Mouredev/mi_fichero.json") as json_file:
    for line in json_file.readlines():
        print(line)

# Carga el fichero json y lo convierte en un diccionario
json_dict = json.load(open("Aprendizaje/Mouredev/mi_fichero.json")) #Carga el fichero json y lo convierte en un diccionario
print(json_dict) # Para ver el diccionario
print(type(json_dict)) # Para ver el tipo de dato en este caso es un diccionario
print(json_dict["nombre"]) # Para acceder a un valor del diccionario
'''

# .csv file
# Es un fichero de texto plano que contiene datos separados por comas

import csv
'''
csv_file = open("Aprendizaje/Mouredev/mi_fichero.csv", "w+") #Crea el fichero csv
csv_writer = csv.writer(csv_file) #Crea el escritor del fichero csv
csv_writer.writerow(["nombre", "apellido", "edad", "lenguaje"]) #Escribe una fila en el fichero csv (Es cómo un array)
csv_writer.writerow(["Xavi", "Quesada", 31, "Python"]) #Escribe una fila en el fichero csv (Es cómo un array)
csv_writer.writerow(["Juan", "", 25, ""])
'''

# .xlsx file
# Es un fichero de Excel

#import xlrd # Debe instalarse la librería xlrd para poder leer ficheros de Excel

# .xml file
# Es un fichero de texto plano que contiene datos separados por etiquetas
'''
import xml

xml_file = open("Aprendizaje/Mouredev/mi_fichero.xml", "w+") #Crea el fichero xml
xml_file.write("<nombre>Xavi</nombre>") #Escribe en el fichero xml
xml_file.write("<apellido>Quesada</apellido>") #Escribe en el fichero xml
xml_file.write("<edad>31</edad>") #Escribe en el fichero xml
xml_file.write("<lenguaje>Python</lenguaje>") #Escribe en el fichero xml
'''

# .html file
# Es un fichero de texto plano que contiene datos separados por etiquetas

html_file = open("Aprendizaje/Mouredev/mi_fichero.html", "w+") #Crea el fichero html
html_file.write("<html><head><title>Mi página web</title></head><body><h1>Hola mundo</h1></body></html>") #Escribe en el fichero html


# .pdf file
# Es un fichero de PDF
'''
from fpdf import FPDF

pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", size=12)
pdf.cell(200, 10, txt="Hola mundo", ln=1, align="C")
pdf.output("Aprendizaje/Mouredev/mi_fichero.pdf")
'''

# .docx file
# Es un fichero de Word
'''
from docx import Document

document = Document()
document.add_heading("Mi documento Word", 0)
document.add_paragraph("Hola mundo")
document.save("Aprendizaje/Mouredev/mi_fichero.docx")
'''

# .pptx file
# Es un fichero de PowerPoint
'''
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Mi presentación"
subtitle.text = "Hola mundo"
prs.save("Aprendizaje/Mouredev/mi_fichero.pptx")
'''

# .mp3 file
# Es un fichero de audio
'''
import os as os   # Debe instalarse la librería os para poder reproducir ficheros de audio

os.system("Aprendizaje/Mouredev/mi_fichero.mp3")
'''

# .mp4 file
# Es un fichero de video
'''
import os as os   # Debe instalarse la librería os para poder reproducir ficheros de video

os.system("Aprendizaje/Mouredev/mi_fichero.mp4")
'''